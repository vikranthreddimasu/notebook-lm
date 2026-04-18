from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable


@dataclass
class LoadedDocument:
    path: Path
    text: str
    metadata: dict[str, str] | None = None


class DocumentLoaderError(Exception):
    """Raised when document loading fails."""


def iter_supported_files(path: Path, recursive: bool = True) -> Iterable[Path]:
    """
    Iterate over supported document files in the given path.
    
    Supported formats: .txt, .md, .pdf, .docx, .pptx, .py
    """
    if path.is_file():
        if _is_supported(path):
            yield path
        return
    
    if recursive:
        pattern = "**/*"
    else:
        pattern = "*"
    
    for file_path in path.glob(pattern):
        if file_path.is_file() and _is_supported(file_path):
            yield file_path


def _is_supported(path: Path) -> bool:
    """Check if file extension is supported."""
    supported = {".txt", ".md", ".pdf", ".docx", ".pptx", ".py"}
    return path.suffix.lower() in supported


def load_document(file_path: Path) -> LoadedDocument:
    """
    Load a document from file path.
    
    Supports: .txt, .md, .pdf, .docx, .pptx, .py
    """
    if not file_path.exists():
        raise DocumentLoaderError(f"File does not exist: {file_path}")
    
    suffix = file_path.suffix.lower()
    
    if suffix in {".txt", ".md"}:
        return _load_text_file(file_path)
    elif suffix == ".pdf":
        return _load_pdf(file_path)
    elif suffix == ".docx":
        return _load_docx(file_path)
    elif suffix == ".pptx":
        return _load_pptx(file_path)
    elif suffix == ".py":
        return _load_text_file(file_path)
    else:
        raise DocumentLoaderError(f"Unsupported file type: {suffix}")


def _load_text_file(file_path: Path) -> LoadedDocument:
    """Load plain text or markdown file."""
    try:
        text = file_path.read_text(encoding="utf-8")
        return LoadedDocument(path=file_path, text=text)
    except Exception as e:
        raise DocumentLoaderError(f"Failed to read {file_path}: {e}")


def _load_pdf(file_path: Path) -> LoadedDocument:
    """Load PDF file using pypdf."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        text_parts = []
        page_count = 0

        for page in reader.pages:
            page_count += 1
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

        text = "\n\n".join(part for part in text_parts if part)

        # Scanned PDFs (image-only, no embedded text) produce near-empty
        # output. Previously we silently "succeeded" with 0 chunks and the
        # user saw "no relevant documents found" on every query. Raise a
        # clear error instead; the API turns this into a 400 with an
        # actionable message.
        MIN_CHARS_PER_PAGE = 40
        expected_min = MIN_CHARS_PER_PAGE * max(1, page_count)
        if len(text.strip()) < min(100, expected_min):
            raise DocumentLoaderError(
                "This PDF appears to be a scan (no embedded text). OCR isn't "
                "supported yet — run the file through a tool like Adobe or "
                "Preview's OCR first, then re-upload."
            )

        return LoadedDocument(path=file_path, text=text)
    except ImportError:
        raise DocumentLoaderError("pypdf is required for PDF support")
    except DocumentLoaderError:
        raise
    except Exception as e:
        raise DocumentLoaderError(f"Failed to read PDF {file_path}: {e}")


def _load_docx(file_path: Path) -> LoadedDocument:
    """Load DOCX file using python-docx."""
    try:
        from docx import Document
        
        doc = Document(str(file_path))
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        text = "\n\n".join(text_parts)
        return LoadedDocument(path=file_path, text=text)
    except ImportError:
        raise DocumentLoaderError("python-docx is required for DOCX support")
    except Exception as e:
        raise DocumentLoaderError(f"Failed to read DOCX {file_path}: {e}")


def _load_pptx(file_path: Path) -> LoadedDocument:
    """Load PPTX file using python-pptx by concatenating text from slides."""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        text_parts: list[str] = []
        for slide in prs.slides:
            # Extract text from all shapes that contain text frames
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
                elif hasattr(shape, "text_frame") and shape.text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text:
                            text_parts.append(p.text)
        text = "\n\n".join(part.strip() for part in text_parts if part and part.strip())
        return LoadedDocument(path=file_path, text=text)
    except ImportError:
        raise DocumentLoaderError("python-pptx is required for PPTX support")
    except Exception as e:
        raise DocumentLoaderError(f"Failed to read PPTX {file_path}: {e}")
